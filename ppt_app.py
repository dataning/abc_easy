import streamlit as st
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pydantic import BaseModel
import json
import io
import re
import base64
from PIL import Image
from io import BytesIO
import os
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.patches import Rectangle, FancyBboxPatch
import numpy as np
from matplotlib.patches import Polygon

# Pydantic models for structured output
class SlideContent(BaseModel):
    title: str
    content: str
    slide_type: str = "bullet"  # bullet, title, image, blank
    image_prompt: str = None

class PresentationStructure(BaseModel):
    slides: list[SlideContent]

# Initialize session state
if 'presentation' not in st.session_state:
    st.session_state.presentation = None
if 'chat_history' not in st.session_state:
    st.session_state.chat_history = []
if 'slide_content' not in st.session_state:
    st.session_state.slide_content = []
if 'generated_images' not in st.session_state:
    st.session_state.generated_images = {}
if 'pending_prompt' not in st.session_state:
    st.session_state.pending_prompt = None
if 'selected_template' not in st.session_state:
    st.session_state.selected_template = "Default"

# Template configurations
TEMPLATES = {
    "Default": {
        "bg_color": "#FFFFFF",
        "title_color": "#1F2937",
        "text_color": "#374151",
        "accent_color": "#3B82F6",
        "font_title": "Arial",
        "font_body": "Arial"
    },
    "Corporate": {
        "bg_color": "#F8FAFC",
        "title_color": "#0F172A",
        "text_color": "#475569",
        "accent_color": "#0EA5E9",
        "font_title": "Arial",
        "font_body": "Arial"
    },
    "Creative": {
        "bg_color": "#FEF3C7",
        "title_color": "#92400E",
        "text_color": "#78350F",
        "accent_color": "#F59E0B",
        "font_title": "Arial",
        "font_body": "Arial"
    },
    "Dark": {
        "bg_color": "#1E293B",
        "title_color": "#F1F5F9",
        "text_color": "#CBD5E1",
        "accent_color": "#A78BFA",
        "font_title": "Arial",
        "font_body": "Arial"
    }
}

# Page config
st.set_page_config(page_title="PowerPoint Chat Generator", page_icon="📊", layout="wide")

# Sidebar for API configuration
with st.sidebar:
    st.header("Configuration")
    api_key = st.text_input("Gemini API Key", type="password", help="Get your API key from Google AI Studio")
    
    st.divider()
    st.header("Model Settings")
    text_model = st.selectbox("Text Model", ["gemini-2.0-flash", "gemini-2.5-flash"])
    enable_images = st.checkbox("Enable Image Generation", value=False)
    if enable_images:
        st.info("Image generation uses Imagen 3.0")
    
    st.divider()
    st.header("Presentation Settings")
    template = st.selectbox("Template", list(TEMPLATES.keys()), index=0)
    if template != st.session_state.selected_template:
        st.session_state.selected_template = template
    
    # Show template preview
    st.markdown("**Template Preview:**")
    template_config = TEMPLATES[template]
    col1, col2 = st.columns(2)
    with col1:
        st.color_picker("Background", template_config["bg_color"], disabled=True)
        st.color_picker("Title", template_config["title_color"], disabled=True)
    with col2:
        st.color_picker("Text", template_config["text_color"], disabled=True)
        st.color_picker("Accent", template_config["accent_color"], disabled=True)
    
    if st.button("New Presentation", type="primary"):
        st.session_state.presentation = None
        st.session_state.slide_content = []
        st.session_state.chat_history = []
        st.session_state.generated_images = {}
        st.rerun()

# Main title
st.title("📊 PowerPoint Chat Generator with Gemini")
st.markdown("Generate and modify PowerPoint presentations using Google's Gemini API")

# Initialize Gemini client
def get_gemini_client(api_key):
    if not api_key:
        return None
    return OpenAI(
        api_key=api_key,
        base_url="https://generativelanguage.googleapis.com/v1beta/openai/"
    )

# Helper functions
def create_new_presentation():
    prs = Presentation()
    # Set slide size to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    
    if subtitle and len(slide.placeholders) > 1:
        subtitle_placeholder = slide.placeholders[1]
        subtitle_placeholder.text = subtitle
    
    return slide

def add_content_slide(prs, title, content, layout_type="bullet"):
    if layout_type == "bullet":
        slide_layout = prs.slide_layouts[1]  # Bullet slide layout
    else:
        slide_layout = prs.slide_layouts[5]  # Blank layout
    
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Add content
    if layout_type == "bullet" and len(slide.placeholders) > 1:
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        
        # Parse content into bullet points
        lines = content.strip().split('\n')
        for i, line in enumerate(lines):
            if i == 0:
                text_frame.text = line.strip('- •*')
            else:
                p = text_frame.add_paragraph()
                p.text = line.strip('- •*')
                p.level = 0
    
    return slide

def add_image_slide(prs, title, image_path_or_data, content=""):
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_left = Inches(0.5)
    title_top = Inches(0.5)
    title_width = Inches(12.333)
    title_height = Inches(1)
    
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    
    # Add image
    img_left = Inches(1.5)
    img_top = Inches(2)
    img_width = Inches(10.333)
    
    if isinstance(image_path_or_data, bytes):
        # If image data is bytes, save to temporary file
        temp_path = "temp_image.png"
        img = Image.open(BytesIO(image_path_or_data))
        img.save(temp_path)
        slide.shapes.add_picture(temp_path, img_left, img_top, width=img_width)
        os.remove(temp_path)
    else:
        slide.shapes.add_picture(image_path_or_data, img_left, img_top, width=img_width)
    
    # Add content if provided
    if content:
        content_left = Inches(0.5)
        content_top = Inches(6)
        content_width = Inches(12.333)
        content_height = Inches(1)
        
        content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
        content_frame = content_box.text_frame
        content_frame.text = content
        content_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def generate_image(client, prompt):
    """Generate image using Imagen 3.0"""
    try:
        response = client.images.generate(
            model="imagen-3.0-generate-002",
            prompt=prompt,
            response_format='b64_json',
            n=1,
        )
        
        image_data = base64.b64decode(response.data[0].b64_json)
        return image_data
    except Exception as e:
        st.error(f"Image generation error: {str(e)}")
        return None

def generate_presentation_content_structured(prompt, client, model):
    """Generate presentation content using structured output"""
    system_prompt = """You are a PowerPoint presentation expert. When asked to create or modify presentations, 
    respond with a well-structured presentation. For each slide, provide:
    - A clear, concise title
    - Bullet points or content (use • or - for bullets)
    - Specify slide_type: "title" for title slides, "bullet" for bullet points, "image" for slides that should have images
    - For image slides, provide an image_prompt describing what image should be generated
    
    Make presentations engaging, professional, and well-organized."""
    
    try:
        completion = client.beta.chat.completions.parse(
            model=model,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt}
            ],
            response_format=PresentationStructure,
            temperature=0.7
        )
        
        return completion.choices[0].message.parsed, None
    except Exception as e:
        # Fallback to regular completion
        try:
            response = client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7
            )
            return parse_text_response(response.choices[0].message.content), None
        except Exception as e2:
            return None, str(e2)

def parse_text_response(response):
    """Parse text response into PresentationStructure"""
    slides = []
    
    # Split by "Slide" markers
    slide_sections = re.split(r'Slide \d+:', response)
    
    for section in slide_sections[1:]:  # Skip first empty section
        lines = section.strip().split('\n')
        if lines:
            title = lines[0].strip()
            content = '\n'.join(lines[1:]).strip()
            
            # Detect if this should be an image slide
            slide_type = "bullet"
            image_prompt = None
            
            if "image:" in content.lower() or "[image" in content.lower():
                slide_type = "image"
                # Extract image description
                image_match = re.search(r'\[image:?\s*(.*?)\]', content, re.IGNORECASE)
                if image_match:
                    image_prompt = image_match.group(1)
                    content = re.sub(r'\[image:?\s*.*?\]', '', content, flags=re.IGNORECASE).strip()
            
            slides.append(SlideContent(
                title=title,
                content=content,
                slide_type=slide_type,
                image_prompt=image_prompt
            ))
    
    return PresentationStructure(slides=slides)

def create_slide_preview(slide, slide_index, template_config):
    """Create a consistent visual preview of a slide using matplotlib"""
    fig, ax = plt.subplots(1, 1, figsize=(10, 5.625), dpi=100)
    
    # Convert hex colors to RGB for matplotlib
    bg_rgb = np.array(hex_to_rgb(template_config["bg_color"])) / 255
    title_rgb = np.array(hex_to_rgb(template_config["title_color"])) / 255
    text_rgb = np.array(hex_to_rgb(template_config["text_color"])) / 255
    accent_rgb = np.array(hex_to_rgb(template_config["accent_color"])) / 255
    
    # Set figure background
    fig.patch.set_facecolor(bg_rgb)
    ax.set_facecolor(bg_rgb)
    
    # Add subtle border
    border = Rectangle((0.05, 0.05), 15.9, 8.9, 
                      fill=False, 
                      edgecolor='gray', 
                      linewidth=1.5,
                      alpha=0.3)
    ax.add_patch(border)
    
    # Determine slide type and layout
    slide_type = slide.get('slide_type', 'bullet')
    
    if slide_type == 'title':
        # Title slide layout
        # Add decorative element
        accent_rect = Rectangle((0, 7), 16, 2, 
                               facecolor=accent_rgb, 
                               alpha=0.1)
        ax.add_patch(accent_rect)
        
        # Center title
        ax.text(8, 5.5, slide['title'], 
                fontsize=28, 
                weight='bold', 
                ha='center', 
                va='center',
                color=title_rgb,
                wrap=True)
        
        # Subtitle/content
        if slide.get('content'):
            ax.text(8, 3.5, slide['content'][:100], 
                    fontsize=16, 
                    ha='center', 
                    va='center',
                    color=text_rgb,
                    style='italic',
                    wrap=True)
    
    elif slide_type == 'image':
        # Image slide layout
        # Title area with accent
        title_bg = Rectangle((0, 7.5), 16, 1.5, 
                           facecolor=accent_rgb, 
                           alpha=0.05)
        ax.add_patch(title_bg)
        
        # Title
        ax.text(1, 8.2, slide['title'], 
                fontsize=20, 
                weight='bold', 
                ha='left', 
                va='center',
                color=title_rgb)
        
        # Image placeholder
        if slide_index in st.session_state.generated_images:
            # Show actual image thumbnail
            img_data = st.session_state.generated_images[slide_index]
            img = Image.open(BytesIO(img_data))
            
            # Calculate position to center image
            img_width = 10
            img_height = 5
            img_x = (16 - img_width) / 2
            img_y = 1.5
            
            # Add image
            ax.imshow(img, extent=[img_x, img_x + img_width, img_y, img_y + img_height], aspect='auto')
            
            # Add subtle frame
            img_frame = Rectangle((img_x - 0.1, img_y - 0.1), 
                                img_width + 0.2, 
                                img_height + 0.2,
                                fill=False,
                                edgecolor=accent_rgb,
                                linewidth=2)
            ax.add_patch(img_frame)
        else:
            # Placeholder for image
            img_placeholder = FancyBboxPatch((3, 1.5), 10, 5,
                                           boxstyle="round,pad=0.1",
                                           facecolor='lightgray',
                                           edgecolor=accent_rgb,
                                           linewidth=2,
                                           alpha=0.3)
            ax.add_patch(img_placeholder)
            
            # Image icon and text
            ax.text(8, 4, '🖼️', fontsize=40, ha='center', va='center', alpha=0.5)
            if slide.get('image_prompt'):
                ax.text(8, 2.5, f'"{slide["image_prompt"][:50]}..."', 
                       fontsize=12, 
                       ha='center', 
                       va='center',
                       style='italic',
                       color=text_rgb,
                       alpha=0.7)
    
    else:  # bullet slide
        # Title area with subtle accent
        title_bg = Rectangle((0, 7.5), 16, 1.5, 
                           facecolor=accent_rgb, 
                           alpha=0.05)
        ax.add_patch(title_bg)
        
        # Title with accent line
        ax.text(1, 8.2, slide['title'], 
                fontsize=20, 
                weight='bold', 
                ha='left', 
                va='center',
                color=title_rgb)
        
        # Accent line under title
        ax.plot([1, 15], [7.5, 7.5], color=accent_rgb, linewidth=2, alpha=0.5)
        
        # Content area
        if slide.get('content'):
            lines = [line.strip() for line in slide['content'].split('\n') if line.strip()][:6]
            y_pos = 6.5
            
            for i, line in enumerate(lines):
                # Bullet point
                ax.plot(1.5, y_pos, 'o', color=accent_rgb, markersize=6)
                
                # Text (truncate if too long)
                text = line.strip('- •*')
                if len(text) > 80:
                    text = text[:77] + '...'
                
                ax.text(2, y_pos, text, 
                       fontsize=14, 
                       va='center',
                       color=text_rgb)
                
                y_pos -= 1
                
                # Add "..." if more content
                if i == 5 and len(lines) > 6:
                    ax.text(8, y_pos, '...', 
                           fontsize=14, 
                           ha='center',
                           va='center',
                           color=text_rgb,
                           style='italic')
    
    # Add slide number
    slide_num_bg = patches.Circle((14.5, 0.8), 0.4, 
                                 facecolor=accent_rgb, 
                                 alpha=0.2)
    ax.add_patch(slide_num_bg)
    ax.text(14.5, 0.8, str(slide_index + 1), 
           fontsize=10, 
           ha='center', 
           va='center',
           color=accent_rgb,
           weight='bold')
    
    # Set axis properties
    ax.set_xlim(0, 16)
    ax.set_ylim(0, 9)
    ax.set_aspect('equal')
    ax.axis('off')
    
    # Tight layout
    plt.tight_layout(pad=0)
    
    # Convert to image
    buf = io.BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight', dpi=100, facecolor=bg_rgb)
    buf.seek(0)
    plt.close()
    
    return buf

def save_presentation(prs):
    """Save presentation to bytes buffer for download"""
    try:
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        return pptx_buffer.getvalue()
    except Exception as e:
        st.error(f"Error saving presentation: {str(e)}")
        return None

# Chat interface
col1, col2 = st.columns([3, 2])

with col1:
    st.header("Chat Interface")
    
    # Display chat history
    chat_container = st.container()
    with chat_container:
        for msg in st.session_state.chat_history:
            with st.chat_message(msg["role"]):
                st.write(msg["content"])
    
    # Check for pending prompt
    if st.session_state.pending_prompt:
        prompt = st.session_state.pending_prompt
        st.session_state.pending_prompt = None
    else:
        # Chat input
        prompt = st.chat_input("Describe the presentation you want to create or modify...")
    
    if prompt:
        # Add user message to chat
        st.session_state.chat_history.append({"role": "user", "content": prompt})
        
        with st.chat_message("user"):
            st.write(prompt)
        
        # Generate response
        with st.chat_message("assistant"):
            client = get_gemini_client(api_key)
            
            if not client:
                st.error("Please provide a Gemini API key")
            else:
                with st.spinner("Generating presentation content..."):
                    result, error = generate_presentation_content_structured(prompt, client, text_model)
                    
                    if error:
                        st.error(f"Error: {error}")
                    else:
                        # Display the generated structure
                        response_text = f"I'll create a presentation with {len(result.slides)} slides:\n\n"
                        for i, slide in enumerate(result.slides):
                            response_text += f"**Slide {i+1}: {slide.title}**\n"
                            if slide.slide_type == "image" and slide.image_prompt:
                                response_text += f"*[Image: {slide.image_prompt}]*\n"
                            if slide.content:
                                response_text += f"{slide.content}\n"
                            response_text += "\n"
                        
                        st.write(response_text)
                        st.session_state.chat_history.append({"role": "assistant", "content": response_text})
                        
                        # Create slides
                        if st.session_state.presentation is None:
                            st.session_state.presentation = create_new_presentation()
                        
                        progress = st.progress(0)
                        for i, slide_data in enumerate(result.slides):
                            progress.progress((i + 1) / len(result.slides))
                            
                            if slide_data.slide_type == "title":
                                add_title_slide(
                                    st.session_state.presentation,
                                    slide_data.title,
                                    slide_data.content
                                )
                            elif slide_data.slide_type == "image" and enable_images and slide_data.image_prompt:
                                # Generate image
                                with st.spinner(f"Generating image for slide {i+1}..."):
                                    image_data = generate_image(client, slide_data.image_prompt)
                                    if image_data:
                                        st.session_state.generated_images[i] = image_data
                                        add_image_slide(
                                            st.session_state.presentation,
                                            slide_data.title,
                                            image_data,
                                            slide_data.content
                                        )
                                    else:
                                        # Fallback to bullet slide
                                        add_content_slide(
                                            st.session_state.presentation,
                                            slide_data.title,
                                            slide_data.content,
                                            "bullet"
                                        )
                            else:
                                add_content_slide(
                                    st.session_state.presentation,
                                    slide_data.title,
                                    slide_data.content,
                                    slide_data.slide_type
                                )
                        
                        st.session_state.slide_content.extend([s.dict() for s in result.slides])
                        progress.empty()
                        st.success(f"✅ Added {len(result.slides)} slides to the presentation!")
                        st.rerun()

with col2:
    st.header("Presentation Preview")
    
    if st.session_state.presentation:
        # Show slide count and template
        col_info1, col_info2 = st.columns(2)
        with col_info1:
            st.info(f"📊 Total slides: {len(st.session_state.presentation.slides)}")
        with col_info2:
            st.info(f"🎨 Template: {st.session_state.selected_template}")
        
        # Display slide previews
        if st.session_state.slide_content:
            st.subheader("Slide Previews")
            
            # Get current template config
            template_config = TEMPLATES[st.session_state.selected_template]
            
            # Create tabs for different view modes
            tab1, tab2, tab3 = st.tabs(["📑 All Slides", "🔍 Detailed View", "⚡ Quick View"])
            
            with tab1:
                # Grid view of all slides
                st.markdown("**All slides at a glance:**")
                
                # Calculate grid layout
                num_slides = len(st.session_state.slide_content)
                cols_per_row = 2
                
                for row in range(0, num_slides, cols_per_row):
                    cols = st.columns(cols_per_row)
                    for col_idx in range(cols_per_row):
                        slide_idx = row + col_idx
                        if slide_idx < num_slides:
                            with cols[col_idx]:
                                slide = st.session_state.slide_content[slide_idx]
                                
                                # Generate preview
                                preview = create_slide_preview(slide, slide_idx, template_config)
                                st.image(preview, use_container_width=True)
                                
                                # Slide info
                                st.caption(f"**Slide {slide_idx + 1}:** {slide['title'][:30]}...")
                                
                                # Quick actions as a single row
                                if st.button(f"📝 Edit Slide {slide_idx + 1}", key=f"edit_{slide_idx}", use_container_width=True):
                                    st.session_state.pending_prompt = f"Modify slide {slide_idx + 1}: {slide['title']}"
                                    st.rerun()
                                
                                if slide.get('slide_type') == 'image' and slide_idx in st.session_state.generated_images:
                                    if st.button(f"🖼️ View Image", key=f"view_img_{slide_idx}", use_container_width=True):
                                        st.image(st.session_state.generated_images[slide_idx], caption=f"Generated image for slide {slide_idx + 1}")
            
            with tab2:
                # Detailed view with expandable sections
                for i, slide in enumerate(st.session_state.slide_content):
                    with st.expander(f"Slide {i+1}: {slide['title']}", expanded=(i == 0)):
                        # Preview
                        preview = create_slide_preview(slide, i, template_config)
                        st.image(preview, use_container_width=True)
                        
                        # Details
                        st.markdown("**Content:**")
                        if slide.get('content'):
                            st.text_area("", slide['content'], height=150, disabled=True, key=f"content_{i}")
                        else:
                            st.info("No content")
                        
                        st.markdown("**Properties:**")
                        st.write(f"Type: `{slide.get('slide_type', 'bullet')}`")
                        if slide.get('image_prompt'):
                            st.write(f"Image prompt: `{slide['image_prompt'][:50]}...`")
                        
                        if i in st.session_state.generated_images:
                            st.markdown("**Generated Image:**")
                            st.image(st.session_state.generated_images[i], use_container_width=True)
                        
                        # Edit button
                        if st.button(f"✏️ Edit This Slide", key=f"edit_detail_{i}", use_container_width=True):
                            st.session_state.pending_prompt = f"Modify slide {i + 1}: {slide['title']}"
                            st.rerun()
            
            with tab3:
                # Quick text view
                st.markdown("**Quick outline view:**")
                for i, slide in enumerate(st.session_state.slide_content):
                    st.markdown(f"**{i+1}. {slide['title']}**")
                    if slide.get('slide_type') == 'image' and slide.get('image_prompt'):
                        st.markdown(f"   - *Image: {slide['image_prompt']}*")
                    if slide.get('content'):
                        lines = slide['content'].split('\n')[:3]
                        for line in lines:
                            if line.strip():
                                st.markdown(f"   - {line.strip('- •*')[:50]}...")
                    st.markdown("")
        
        # Download section
        st.divider()
        pptx_data = save_presentation(st.session_state.presentation)
        if pptx_data:
            st.download_button(
                label="📥 Download PowerPoint Presentation",
                data=pptx_data,
                file_name="gemini_generated_presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                type="primary"
            )
        
        # Modification suggestions
        st.divider()
        st.subheader("Quick Actions")
        col_sug1, col_sug2 = st.columns(2)
        
        suggestions_left = [
            "Add a slide about implementation timeline",
            "Include a summary slide with key takeaways",
            "Add a Q&A slide at the end"
        ]
        
        suggestions_right = [
            "Add more visual elements",
            "Create a comparison table",
            "Add team introduction slide"
        ]
        
        with col_sug1:
            for i, suggestion in enumerate(suggestions_left):
                if st.button(suggestion, key=f"mod_left_{i}", use_container_width=True):
                    st.session_state.pending_prompt = suggestion
                    st.rerun()
        
        with col_sug2:
            for i, suggestion in enumerate(suggestions_right):
                if st.button(suggestion, key=f"mod_right_{i}", use_container_width=True):
                    st.session_state.pending_prompt = suggestion
                    st.rerun()
        
        # Custom modification
        st.markdown("**Custom Modification:**")
        custom_mod = st.text_input("Describe your modification...", key="custom_mod")
        if st.button("Apply Custom Modification", use_container_width=True):
            if custom_mod:
                st.session_state.pending_prompt = custom_mod
                st.rerun()
    
    else:
        # Empty state
        st.info("👋 Start by describing the presentation you want to create in the chat!")
        
        # Template showcase
        st.subheader("Available Templates")
        template_cols = st.columns(2)
        for i, (template_name, config) in enumerate(TEMPLATES.items()):
            with template_cols[i % 2]:
                # Create mini preview of template
                fig, ax = plt.subplots(1, 1, figsize=(4, 2.25), dpi=80)
                
                bg_rgb = np.array(hex_to_rgb(config["bg_color"])) / 255
                title_rgb = np.array(hex_to_rgb(config["title_color"])) / 255
                accent_rgb = np.array(hex_to_rgb(config["accent_color"])) / 255
                
                fig.patch.set_facecolor(bg_rgb)
                ax.set_facecolor(bg_rgb)
                
                # Sample layout
                ax.add_patch(Rectangle((0, 3.5), 8, 1, facecolor=accent_rgb, alpha=0.1))
                ax.text(4, 4, template_name, fontsize=12, weight='bold', ha='center', va='center', color=title_rgb)
                
                # Sample bullets
                for j in range(3):
                    ax.plot(0.5, 2.5 - j*0.5, 'o', color=accent_rgb, markersize=4)
                    ax.plot([0.8, 7.5], [2.5 - j*0.5, 2.5 - j*0.5], color=title_rgb, alpha=0.3, linewidth=8)
                
                ax.set_xlim(0, 8)
                ax.set_ylim(0, 4.5)
                ax.axis('off')
                
                buf = io.BytesIO()
                plt.savefig(buf, format='png', bbox_inches='tight', dpi=80, facecolor=bg_rgb)
                buf.seek(0)
                plt.close()
                
                st.image(buf, use_container_width=True)
                st.caption(f"**{template_name}** Template")
        
        st.divider()
        st.subheader("Example Prompts")
        examples = [
            "Create a 5-slide presentation about the future of AI",
            "Make a startup pitch deck for an eco-friendly product",
            "Generate a technical presentation about cloud computing",
            "Create a training deck on Python best practices",
            "Build a sales presentation for a SaaS product",
            "Design a company overview presentation"
        ]
        
        example_cols = st.columns(2)
        for i, example in enumerate(examples):
            with example_cols[i % 2]:
                if st.button(example, key=f"ex_{i}", use_container_width=True):
                    st.session_state.pending_prompt = example
                    st.rerun()

# Footer with instructions
with st.expander("ℹ️ How to use this app", expanded=False):
    col_inst1, col_inst2 = st.columns(2)
    
    with col_inst1:
        st.markdown("""
        ### 🚀 Quick Start
        1. **Get API Key**: Visit [Google AI Studio](https://aistudio.google.com/apikey)
        2. **Enter API Key**: Paste it in the sidebar
        3. **Choose Template**: Select a design theme
        4. **Describe Your Presentation**: Use natural language
        5. **Review & Download**: Get your .pptx file
        
        ### ✨ Features
        - **Consistent Previews**: See exactly how slides will look
        - **Multiple Templates**: Choose from different design themes
        - **Smart Layouts**: Automatic formatting for different content types
        - **Image Generation**: AI-powered images with Imagen 3.0
        """)
    
    with col_inst2:
        st.markdown("""
        ### 💡 Pro Tips
        - Be specific about the number of slides and topics
        - Mention if you want images on certain slides
        - Use modification suggestions for quick edits
        - Try different templates for different moods
        - Enable image generation for visual presentations
        
        ### 📝 Example Formats
        - **Bullet slides**: Great for key points
        - **Title slides**: Perfect for section breaks
        - **Image slides**: Ideal for visual impact
        - **Mixed content**: Combine text and visuals
        """)

# Footer
st.divider()
col1, col2, col3 = st.columns(3)
with col1:
    st.caption("🚀 Powered by Google Gemini & Imagen")
with col2:
    st.caption("📊 Built with Streamlit & python-pptx")
with col3:
    st.caption("🎨 Consistent slide previews")